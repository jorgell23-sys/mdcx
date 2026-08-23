# Copyright 2026 Jorge Ellena G.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Reference text extraction from source documents.

Provides the text a document actually exposes, read with a library
independent from the engine that performed the conversion, so that fidelity
is measured against an external reference rather than against itself.
"""
from __future__ import annotations

from html.parser import HTMLParser
from pathlib import Path


# What share of a document's pages must carry no text before it is treated as
# a scan. A document below this is one with some blank or pictorial pages, which
# is every book; a document above it is one that was photographed.
OCR_SHARE = 0.5


def _pdf_text(path: Path) -> tuple[str, dict]:
    from . import pdf as _pdf

    doc = _pdf.open_document(path)
    parts: list[str] = []
    pages_without_text = 0
    images = 0
    try:
        for page in doc:
            txt = _pdf.page_text(page)
            if len(txt.strip()) < 20:
                pages_without_text += 1
            images += _pdf.count_images(page)
            parts.append(txt)
        # Optical recognition is applied to the whole document, so it has to be
        # decided about the whole document. A single page without text is
        # ordinary -- a half-title, a blank verso, a page holding one figure --
        # and a book of six hundred pages will nearly always have one. Asking
        # for OCR on that basis sends the entire document down the most
        # expensive path there is, to recover a page that has nothing on it.
        # Measured on a real run: nine chapters in a hundred and seventy-seven
        # were converted that way, each of them a document whose text the cheap
        # extractor was already reading in full.
        #
        # What OCR is for is a document that was scanned, and there most of the
        # pages carry no text rather than one of them.
        empty = pages_without_text / len(doc) if len(doc) else 0.0
        meta = {
            "pages": len(doc),
            "pages_without_text": pages_without_text,
            "pages_without_text_share": round(empty, 4),
            "embedded_images": images,
            "needs_ocr": empty >= OCR_SHARE,
        }
    finally:
        doc.close()
    return "\n".join(parts), meta


def _docx_text(path: Path) -> tuple[str, dict]:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    tables = 0
    for table in d.tables:
        tables += 1
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    # Headers and footers usually carry the document code and revision: keep them.
    for section in d.sections:
        for container in (section.header, section.footer):
            for p in container.paragraphs:
                parts.append(p.text)
    return "\n".join(parts), {"tables": tables, "paragraphs": len(d.paragraphs)}


def _xlsx_text(path: Path) -> tuple[str, dict]:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []
    sheets: list[str] = []
    cells = 0
    try:
        for ws in wb.worksheets:
            sheets.append(ws.title)
            parts.append(str(ws.title))
            for row in ws.iter_rows(values_only=True):
                for value in row:
                    if value is None:
                        continue
                    cells += 1
                    parts.append(str(value))
    finally:
        wb.close()
    return "\n".join(parts), {"sheets": sheets, "sheet_count": len(sheets), "cells": cells}


def _pptx_text(path: Path) -> tuple[str, dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts), {"slides": len(prs.slides.__iter__.__self__._sldIdLst)}


def _plain_text(path: Path) -> tuple[str, dict]:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            text = path.read_text(encoding=encoding)
            return text, {"encoding": encoding, "lines": text.count("\n") + 1}
        except UnicodeDecodeError:
            continue
    text = path.read_text(encoding="utf-8", errors="replace")
    return text, {"encoding": "utf-8/replace", "lines": text.count("\n") + 1}



class _TextCollector(HTMLParser):
    """Collects the text of an HTML document, leaving out its markup.

    Script and style hold code rather than prose, and counting it as text would
    credit a conversion for content no reader ever sees.
    """

    _SKIP = {"script", "style", "head", "title", "meta", "link"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self._SKIP:
            self._depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP and self._depth:
            self._depth -= 1

    def handle_data(self, data: str) -> None:
        if not self._depth and data.strip():
            self.parts.append(data)


def _html_to_text(markup: str) -> str:
    collector = _TextCollector()
    try:
        collector.feed(markup)
    except Exception:  # noqa: BLE001 - malformed markup yields what was parsed
        pass
    return " ".join(collector.parts)


def _epub_text(path: Path) -> tuple[str, dict]:
    """Reference text of an EPUB, read from the documents it declares.

    The reading order comes from the spine when the package file can be read.
    Order does not affect coverage, which compares multisets of tokens, but it
    makes the reference readable when someone inspects it.
    """
    import re
    import zipfile

    parts: list[str] = []
    documents = 0
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        content = [n for n in names
                     if n.lower().endswith((".xhtml", ".html", ".htm"))]

        # The package file lists the documents in reading order. It is read as
        # a hint: a spine that names files the archive does not hold is a
        # damaged EPUB, and the remaining documents are still worth reading.
        opf = next((n for n in names if n.lower().endswith(".opf")), None)
        if opf:
            try:
                package = archive.read(opf).decode("utf-8", "replace")
                base = opf.rsplit("/", 1)[0] + "/" if "/" in opf else ""
                paths = dict(re.findall(r'id="([^"]+)"[^>]*href="([^"]+)"', package))
                order = [base + paths[i] for i in re.findall(r'idref="([^"]+)"', package)
                         if i in paths]
                ordered = [n for n in order if n in content]
                content = ordered + [n for n in content if n not in ordered]
            except Exception:  # noqa: BLE001 - the spine is a hint, not a requirement
                pass

        for name in content:
            try:
                text = _html_to_text(archive.read(name).decode("utf-8", "replace"))
            except Exception:  # noqa: BLE001 - a damaged member is skipped, not fatal
                continue
            if text.strip():
                parts.append(text)
                documents += 1

    text = "\n\n".join(parts)
    return text, {"documents": documents, "characters": len(text)}


_EXTRACTORS = {
    "pdf": _pdf_text,
    "docx": _docx_text,
    "xlsx": _xlsx_text,
    "pptx": _pptx_text,
    "epub": _epub_text,
    "text": _plain_text,
    "html": _plain_text,
}


def reference_text(path: Path, kind: str) -> tuple[str, dict]:
    """Return (raw_text, metadata) from the original. Never raises: errors are reported."""
    fn = _EXTRACTORS.get(kind)
    if fn is None:
        return "", {"error": f"no reference extractor for kind={kind}"}
    try:
        return fn(path)
    except Exception as exc:  # noqa: BLE001 - the failure is reported, not raised
        return "", {"error": f"{type(exc).__name__}: {exc}"}
