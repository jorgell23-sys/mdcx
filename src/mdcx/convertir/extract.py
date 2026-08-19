# Copyright 2026 Jorge Ellena G.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Extraccion de texto de referencia ("ground truth") independiente del motor de conversion.

La idea: el conversor principal (Docling) produce Markdown estructurado, pero puede
perder contenido en tablas complejas, cajetines de planos o celdas ocultas. Para poder
afirmar "sin perdida de informacion" hace falta una segunda lectura, hecha con una
libreria distinta y lo mas cruda posible, contra la cual comparar. Eso es este modulo.
"""
from __future__ import annotations

from pathlib import Path


def _pdf_text(path: Path) -> tuple[str, dict]:
    from . import pdf as _pdf

    doc = _pdf.leer_documento(path)
    parts: list[str] = []
    pages_without_text = 0
    images = 0
    try:
        for page in doc:
            # Texto plano en el orden que declara el propio PDF. En los cajetines de
            # planos aparece igual aunque no forme parrafos, que es lo que se quiere de
            # una referencia: lo que el archivo contiene, sin interpretar.
            txt = _pdf.texto_de_pagina(page)
            if len(txt.strip()) < 20:
                pages_without_text += 1
            images += _pdf.contar_imagenes(page)
            parts.append(txt)
        meta = {
            "pages": len(doc),
            "pages_without_text": pages_without_text,
            "embedded_images": images,
            "needs_ocr": pages_without_text > 0,
        }
    finally:
        doc.close()
    return "\n".join(parts), meta


def _docx_text(path: Path) -> tuple[str, dict]:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    tables = 0
    for table in d.tables:
        tables += 1
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    # Encabezados y pies suelen contener codigo de documento y revision: no perderlos.
    for section in d.sections:
        for container in (section.header, section.footer):
            for p in container.paragraphs:
                parts.append(p.text)
    return "\n".join(parts), {"tables": tables, "paragraphs": len(d.paragraphs)}


def _xlsx_text(path: Path) -> tuple[str, dict]:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []
    sheets: list[str] = []
    cells = 0
    try:
        for ws in wb.worksheets:
            sheets.append(ws.title)
            parts.append(str(ws.title))
            for row in ws.iter_rows(values_only=True):
                for value in row:
                    if value is None:
                        continue
                    cells += 1
                    parts.append(str(value))
    finally:
        wb.close()
    return "\n".join(parts), {"sheets": sheets, "sheet_count": len(sheets), "cells": cells}


def _pptx_text(path: Path) -> tuple[str, dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts), {"slides": len(prs.slides.__iter__.__self__._sldIdLst)}


def _plain_text(path: Path) -> tuple[str, dict]:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            text = path.read_text(encoding=encoding)
            return text, {"encoding": encoding, "lines": text.count("\n") + 1}
        except UnicodeDecodeError:
            continue
    text = path.read_text(encoding="utf-8", errors="replace")
    return text, {"encoding": "utf-8/replace", "lines": text.count("\n") + 1}


_EXTRACTORS = {
    "pdf": _pdf_text,
    "docx": _docx_text,
    "xlsx": _xlsx_text,
    "pptx": _pptx_text,
    "text": _plain_text,
    "html": _plain_text,
}


def reference_text(path: Path, kind: str) -> tuple[str, dict]:
    """Devuelve (texto_crudo, metadatos) del original. Nunca lanza: informa el error."""
    fn = _EXTRACTORS.get(kind)
    if fn is None:
        return "", {"error": f"sin extractor de referencia para kind={kind}"}
    try:
        return fn(path)
    except Exception as exc:  # noqa: BLE001 - el fallo se reporta, no interrumpe el lote
        return "", {"error": f"{type(exc).__name__}: {exc}"}
